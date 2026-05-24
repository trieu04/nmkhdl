"""Trích text từ pptx ra markdown đơn giản để đọc."""
import sys, os
from pathlib import Path
from pptx import Presentation

def extract(path):
    p = Presentation(path)
    out = [f"# {Path(path).stem}\n"]
    for i, slide in enumerate(p.slides, 1):
        out.append(f"\n## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        out.append(f"- {t}")
            elif shape.shape_type == 19:  # table
                pass
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n"," ") for c in row.cells]
                    out.append("| " + " | ".join(cells) + " |")
    return "\n".join(out)

if __name__ == "__main__":
    src = Path(sys.argv[1])
    dst = Path(sys.argv[2])
    dst.mkdir(exist_ok=True)
    for f in sorted(src.glob("*.pptx")):
        if f.name.startswith(".~"):
            continue
        md = extract(f)
        (dst / f"{f.stem}.md").write_text(md, encoding="utf-8")
        print(f.name, "->", len(md), "chars")
